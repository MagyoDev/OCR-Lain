from io import BytesIO
from pathlib import Path

from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ocr_lain.config import OCRConfig
from ocr_lain.extractors.base import BaseExtractor
from ocr_lain.models import ExtractedBlock, ExtractResult
from ocr_lain.ocr_engine import OCREngine


class PPTXExtractor(BaseExtractor):

    def __init__(self, config: OCRConfig):
        self.config = config
        self.ocr_engine = OCREngine(config)

    def extract(self, file_path: Path) -> ExtractResult:
        """
        Extrai texto de uma apresentação .pptx.
        """

        blocks: list[ExtractedBlock] = []
        errors: list[str] = []

        try:
            presentation = Presentation(file_path)

            for slide_index, slide in enumerate(presentation.slides, start=1):
                slide_text = self._extract_slide_text(slide)

                if slide_text:
                    blocks.append(
                        ExtractedBlock(
                            source_file=str(file_path),
                            source_type="pptx",
                            location=f"slide:{slide_index}",
                            method="native_text",
                            text=slide_text,
                            metadata={
                                "slide": slide_index,
                            },
                        )
                    )

                table_texts = self._extract_slide_table_texts(slide)

                for table_index, table_text in enumerate(table_texts, start=1):
                    if table_text:
                        blocks.append(
                            ExtractedBlock(
                                source_file=str(file_path),
                                source_type="pptx",
                                location=f"slide:{slide_index}:table:{table_index}",
                                method="native_table_text",
                                text=table_text,
                                metadata={
                                    "slide": slide_index,
                                    "table": table_index,
                                },
                            )
                        )

                if self.config.ocr_embedded_images:
                    image_blocks, image_errors = self._extract_slide_images(
                        file_path=file_path,
                        slide=slide,
                        slide_index=slide_index,
                    )

                    blocks.extend(image_blocks)
                    errors.extend(image_errors)

        except Exception as error:
            errors.append(f"Erro ao processar PPTX {file_path}: {error}")

        return ExtractResult(
            file_path=file_path,
            blocks=blocks,
            errors=errors,
        )

    def _iter_shapes(self, shapes):

        for shape in shapes:
            yield shape

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_shapes(shape.shapes)

    def _extract_slide_text(self, slide) -> str:

        texts: list[str] = []

        for shape in self._iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            if not shape.has_text_frame:
                continue

            paragraphs: list[str] = []

            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()

                if paragraph_text:
                    paragraphs.append(paragraph_text)

            if paragraphs:
                texts.append("\n".join(paragraphs))

        return "\n\n".join(texts)

    def _extract_slide_table_texts(self, slide) -> list[str]:

        table_texts: list[str] = []

        for shape in self._iter_shapes(slide.shapes):
            if not getattr(shape, "has_table", False):
                continue

            if not shape.has_table:
                continue

            rows: list[str] = []

            for row in shape.table.rows:
                cell_values: list[str] = []

                for cell in row.cells:
                    cell_text = cell.text.strip().replace("\n", " ")

                    if cell_text:
                        cell_values.append(cell_text)

                if cell_values:
                    rows.append(" | ".join(cell_values))

            if rows:
                table_texts.append("\n".join(rows))

        return table_texts

    def _extract_slide_images(
        self,
        file_path: Path,
        slide,
        slide_index: int,
    ) -> tuple[list[ExtractedBlock], list[str]]:

        blocks: list[ExtractedBlock] = []
        errors: list[str] = []

        image_index = 0

        for shape in self._iter_shapes(slide.shapes):
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_index += 1

            try:
                image_data = shape.image.blob
                image = Image.open(BytesIO(image_data))

                text = self.ocr_engine.image_to_text(image)

                if text:
                    blocks.append(
                        ExtractedBlock(
                            source_file=str(file_path),
                            source_type="pptx",
                            location=f"slide:{slide_index}:image:{image_index}",
                            method="ocr_embedded_image",
                            text=text,
                            metadata={
                                "slide": slide_index,
                                "image_index": image_index,
                                "format": image.format,
                                "width": image.width,
                                "height": image.height,
                            },
                        )
                    )

            except UnidentifiedImageError:
                errors.append(
                    f"Imagem não reconhecida no slide {slide_index}, imagem {image_index}, arquivo {file_path}"
                )

            except Exception as error:
                errors.append(
                    f"Erro ao aplicar OCR na imagem {image_index} do slide {slide_index} de {file_path}: {error}"
                )

        return blocks, errors